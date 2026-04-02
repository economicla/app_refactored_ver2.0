
"""

DocumentIngestionUseCase - Doküman yükleme işlem hattı (Business Logic)

Ana akış: PDF → Markdown (structured extraction) → Intelligent Chunking → Vector DB + Reranker.
Desteklenen formatlar: PDF (istihbarat raporu: markdown tabanlı extraction), DOCX, TXT, HTML, XLSX, PPTX.

"""

import logging
import re
import time

from datetime import datetime

from pathlib import Path

from typing import List, Optional
 
from app_refactored.core.interfaces import (

    IEmbeddingService,

    IDocumentRepository

)

from app_refactored.core.entities import (

    DocumentChunk,

    DocumentIngestionResult

)

from .document_preprocessing import DocumentPreprocessor
from .intelligent_chunking import IntelligentChunker
from .html_structured_converter import HTMLStructuredConverter

from app_refactored.structured_extractors.vlm_pdf_extractor import VLMPDFExtractor
 
logger = logging.getLogger(__name__)
 
 
class DocumentIngestionUseCase:

    """Doküman yükleme use case - Dependency Injection ile"""

    # HTML format → doküman türü eşlemesi
    DOC_TYPE_MAP = {
        "teklif": "Teklif Özeti",
        "performans": "Performans Değerlendirmesi",
        "mali_veri": "Mali Veri Tabloları",
    }
 
    def __init__(

        self,

        embedding_service: IEmbeddingService,

        document_repository: IDocumentRepository,

        chunk_size: int = 1000,

        chunk_overlap: int = 200,

        vlm_extractor: Optional["VLMPDFExtractor"] = None,

    ):

        """

        Initialize document ingestion use case

        Args:

            embedding_service: IEmbeddingService implementation

            document_repository: IDocumentRepository implementation

            chunk_size: Chunk boyutu (karakterler)

            chunk_overlap: Chunk'lar arası overlap

        """

        self.embedding_service = embedding_service

        self.document_repository = document_repository

        self.chunk_size = chunk_size

        self.chunk_overlap = chunk_overlap

        self.vlm_extractor = vlm_extractor
 
    def _extract_text(self, file_path: str, filename: str) -> str:

        """Dosya türüne göre metni çıkart"""

        path = Path(file_path)

        file_type = path.suffix.lower()

        self._structured_json = None
        self._preprocess_as_markdown = False  # True when text is already markdown (e.g. structured PDF)

        try:

            if file_type == '.pdf':
                if self.vlm_extractor:
                    logger.info("📑 PDF detected → VLM OCR extraction (all PDFs)")
                    return None  # VLM extraction handled async in execute()

                logger.warning("⚠️ VLM extractor unavailable → falling back to generic PDF extraction")
                return self._extract_pdf(file_path)

            elif file_type == '.docx':

                return self._extract_docx(file_path)

            elif file_type == '.txt':

                return self._extract_txt(file_path)

            elif file_type in ('.html', '.htm'):
                return self._extract_html(file_path)

            elif file_type == '.xlsx':

                return self._extract_xlsx(file_path)

            elif file_type == '.pptx':

                return self._extract_pptx(file_path)

            else:

                raise ValueError(f"Unsupported file type: {file_type}")

        except Exception as e:

            logger.error(f"❌ Text extraction failed: {str(e)}")

            raise
 
    def _is_credit_intelligence_pdf(self, file_path: str, filename: str) -> bool:
        """Detect whether this PDF is a Kredi İstihbarat Raporu."""
        fname = filename.lower()
        if "istihbarat" in fname:
            return True

        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:3]:
                    text = (page.extract_text() or "").lower()
                    signals = ("banka istihbarat", "memzuç bilgi", "memzuc bilgi",
                               "limit risk bilgi", "kredi istihbarat")
                    if sum(1 for s in signals if s in text) >= 2:
                        return True
        except Exception as e:
            logger.warning(f"⚠️ PDF intelligence detection failed (treating as generic): {type(e).__name__}: {e}")
        return False


    async def _extract_pdf_vlm(self, file_path: str, original_filename: str) -> str:
        """VLM (Qwen3-VL) ile PDF sayfalarını görüntü olarak okuyup markdown üret."""
        result = await self.vlm_extractor.extract(
            file_path,
            source_display_name=original_filename,
        )

        self._structured_json = None
        errors = result.get("errors", [])
        if errors:
            for e in errors[:5]:
                logger.warning(f"⚠️ VLM page error: {e}")

        text = result["markdown"]
        meta = result["meta"]
        logger.info(
            f"✅ VLM PDF extraction: {len(text)} chars, "
            f"{meta['pages']} pages, {meta['errors']} errors"
        )
        self._preprocess_as_markdown = True
        return text

    def _extract_pdf(self, file_path: str) -> str:

        """PDF'den metin çıkart - pdfplumber (tablo destekli) + PyMuPDF fallback"""

        try:

            import pdfplumber

            pages_text = []

            with pdfplumber.open(file_path) as pdf:

                for page_num, page in enumerate(pdf.pages):

                    page_parts = []

                    page_parts.append(f"## Sayfa {page_num + 1}")

                    # Tabloları çıkar

                    tables = page.extract_tables()

                    if tables:

                        for table in tables:

                            for row in table:

                                cells = [str(c).strip() if c else "" for c in row]

                                non_empty = [c for c in cells if c]
                                #Key:Value çiftlerini tespit et
                                i = 0
                                while i < len(non_empty):
                                    cell = non_empty[i]
                                    if cell.endswith(':') and i + 1 < len(non_empty):
                                        page_parts.append(f"{cell} {non_empty[i+1]}")
                                        i += 2
                                    elif ':' not in cell and i == 0 and len(non_empty) == 2:
                                        page_parts.append(f"{non_empty[0]}: {non_empty[1]}")
                                        i = len(non_empty)
                                    else:
                                        page_parts.append(cell)
                                        i += 1

                                
                    else:

                        text = page.extract_text()

                        if text:

                            page_parts.append(text.strip())

                    if len(page_parts) > 1:

                        pages_text.append("\n".join(page_parts))

            full_text = "\n\n".join(pages_text)

            if len(full_text.strip()) > 100:

                logger.info(f"✅ PDF pdfplumber extraction: {len(full_text)} chars, {len(pages_text)} pages")

                return full_text

            # pdfplumber yetersizse PyMuPDF fallback

            logger.warning(f"⚠️ pdfplumber yetersiz ({len(full_text)} chars), PyMuPDF deneniyor...")

            import fitz

            doc = fitz.open(file_path)

            pages_text = []

            for page_num, page in enumerate(doc):

                text = page.get_text("text")

                if text and text.strip():

                    pages_text.append(f"## Sayfa {page_num + 1}\n\n{text.strip()}")

            doc.close()

            full_text = "\n\n".join(pages_text)

            logger.info(f"✅ PDF PyMuPDF fallback: {len(full_text)} chars")

            return full_text

        except Exception as e:

            logger.error(f"❌ PDF extraction failed: {str(e)}")

            raise
 
    def _extract_docx(self, file_path: str) -> str:

        """DOCX'den metin çıkart"""

        try:

            from docx import Document

            doc = Document(file_path)

            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            logger.info(f"✅ DOCX extraction: {len(text)} chars")

            return text

        except Exception as e:

            logger.error(f"❌ DOCX extraction failed: {str(e)}")

            raise
 
    def _extract_txt(self, file_path: str) -> str:

        """TXT'den metin çıkart"""

        try:

            with open(file_path, 'r', encoding='utf-8') as file:

                text = file.read()

            logger.info(f"✅ TXT extraction: {len(text)} chars")

            return text

        except Exception as e:

            logger.error(f"❌ TXT extraction failed: {str(e)}")

            raise
 
    def _extract_xlsx(self, file_path: str) -> str:

        """XLSX'den metin çıkart"""

        try:

            import openpyxl

            wb = openpyxl.load_workbook(file_path)

            text = ""

            for sheet in wb.sheetnames:

                ws = wb[sheet]

                for row in ws.iter_rows(values_only=True):

                    text += " ".join(str(cell) if cell else "" for cell in row) + "\n"

            logger.info(f"✅ XLSX extraction: {len(text)} chars")

            return text

        except Exception as e:

            logger.error(f"❌ XLSX extraction failed: {str(e)}")

            raise
 
    def _extract_pptx(self, file_path: str) -> str:

        """PPTX'den metin çıkart"""

        try:

            from pptx import Presentation

            prs = Presentation(file_path)

            text = ""

            for slide in prs.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        text += shape.text + "\n"

            logger.info(f"✅ PPTX extraction: {len(text)} chars")

            return text

        except Exception as e:

            logger.error(f"❌ PPTX extraction failed: {str(e)}")

            raise

    def _extract_html(self, file_path: str) -> str:

        """HTML'den metin çıkart - HTMLStructuredConverter ile yapısal parse"""

        try:

            html_content = None

            for enc in ('utf-8', 'windows-1254', 'latin-1', 'iso-8859-9'):

                try:

                    with open(file_path, 'r', encoding=enc) as file:

                        html_content = file.read()

                    logger.info(f"📄 HTML decoded with encoding: {enc}")

                    break

                except UnicodeDecodeError:

                    continue

            if html_content is None:

                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:

                    html_content = file.read()

            converter = HTMLStructuredConverter()

            result = converter.convert(html_content)

            self._detected_html_format = getattr(converter, 'detected_format', None)

            logger.info(f"✅ HTML extraction: {len(result)} chars (format={self._detected_html_format})")

            return result

        except Exception as e:

            logger.error(f"❌ HTML extraction failed: {str(e)}")

            raise

    def _detect_document_type(self, filename: str, text: str = "") -> str:
        """Doküman türünü tespit et (chunk etiketleme için).
        
        HTML: HTMLStructuredConverter formatına göre
        Diğer: dosya adı + içerik analizi
        """
        file_ext = Path(filename).suffix.lower()

        # HTML → converter'ın tespit ettiği formatı kullan
        if file_ext in ('.html', '.htm'):
            html_fmt = getattr(self, '_detected_html_format', None)
            if html_fmt and html_fmt in self.DOC_TYPE_MAP:
                return self.DOC_TYPE_MAP[html_fmt]

        # Dosya adı anahtar kelime tespiti
        fname_lower = filename.lower()
        if "istihbarat" in fname_lower:
            return "İstihbarat Raporu"
        if "teklif" in fname_lower:
            return "Teklif Özeti"
        if "performans" in fname_lower:
            return "Performans Değerlendirmesi"
        if any(kw in fname_lower for kw in ("mali", "solo", "konsolide", "bilanco", "bilanço")):
            return "Mali Veri Tabloları"

        # İçerik analizi (ilk 5000 karakter)
        sample = (text[:5000] if text else "").lower()
        if any(kw in sample for kw in ("istihbarat", "diğer bankalardaki", "bankacılık riskleri")):
            return "İstihbarat Raporu"
        if any(kw in sample for kw in ("limit bilgileri", "teminat koşulları", "rating değerleri")):
            return "Teklif Özeti"
        if any(kw in sample for kw in ("net satışlar", "aktif toplam", "özkaynak")):
            return "Mali Veri Tabloları"

        return "Genel Doküman"

    def _is_dictionary_doc(self, filename: str, text: str = "") -> bool:
        """
        Dosyanın veri sözlüğü dokümanı olup olmadığını tespit et.
        
        Tespit kriterleri:
        1. Dosya adında 'sözlük', 'sozluk', 'dictionary', 'açıklama', 'rehber' geçmesi
        2. İçerikte sözlük kalıpları (alan açıklamaları, tanımlar)
        3. API tarafından is_dictionary=true gönderilmesi (metadata ile)
        """
        fname_lower = filename.lower()
        dict_keywords = ["sözlük", "sozluk", "dictionary", "açıklama", "aciklama",
                         "rehber", "tanım", "tanim", "veri_sozlugu", "data_dict"]
        if any(kw in fname_lower for kw in dict_keywords):
            return True
        
        # İçerik analizi: sözlük dokümanları genellikle alan tanımları içerir
        if text:
            sample = text[:3000].lower()
            definition_patterns = [
                "bu bölüm", "bu alan", "bu sütun", "açıklama:",
                "tanım:", "içerik:", "ne anlama gelir",
                "bu tablo", "bu rapor", "bu doküman",
            ]
            matches = sum(1 for p in definition_patterns if p in sample)
            if matches >= 3:
                return True
        
        return False
 
    
    def _chunk_text(self, text: str) -> List[str]:

        """Metni chunk'lara böl"""

        chunks = []

        start = 0

        while start < len(text):

            end = start + self.chunk_size

            chunk = text[start:end]

            chunks.append(chunk)

            start += (self.chunk_size - self.chunk_overlap)

        logger.info(f"✅ Text chunked: {len(chunks)} chunks")

        return chunks
 
    async def execute(self, file_path: str, filename: str, collection: Optional[str] = None) -> DocumentIngestionResult:

        """

        Execute document ingestion pipeline

        Args:

            file_path: Dosyanın tam yolu

            filename: Orijinal dosya adı

            collection: Dokümanın ait olduğu koleksiyon (ör. 'kredi', 'egitim')

        Returns:

            DocumentIngestionResult

        """

        try:

            pipeline_start = time.monotonic()
            logger.info(f"📥 Ingesting: {filename}")

            # Step 0: Aynı dosyanın eski chunk'larını temizle (re-upload desteği)
            try:
                deleted_count = await self.document_repository.delete_by_filename(filename)
                if deleted_count > 0:
                    logger.info(f"🗑️ Deleted {deleted_count} old chunks for {filename}")
            except Exception as del_err:
                logger.warning(f"⚠️ Old chunk cleanup failed (continuing): {del_err}")

            # Step 1: Metni çıkart
            t0 = time.monotonic()
            logger.info(f"📄 [1/5] Extracting text from {Path(filename).suffix}...")

            text = self._extract_text(file_path, filename)

            if text is None and self.vlm_extractor:
                logger.info(f"📑 Sync extraction returned None → VLM extraction for {Path(filename).name}")
                text = await self._extract_pdf_vlm(file_path, filename)

            if text is None:
                raise ValueError("Text extraction returned None — both sync and VLM paths failed")

            original_length = len(text)
            extract_ms = (time.monotonic() - t0) * 1000
            logger.info(f"📏 [1/5] Extraction done: {original_length:,} chars in {extract_ms:.0f}ms")

            # Step 2: Preprocessing
            t0 = time.monotonic()
            logger.info("🧹 [2/5] Preprocessing text...")
            preprocessor = DocumentPreprocessor()
            raw_type = Path(filename).suffix.lstrip('.')
            preprocess_type = "md" if (
                raw_type in ("html", "htm") or getattr(self, "_preprocess_as_markdown", False)
            ) else raw_type
            text = preprocessor.preprocess(
                text,
                file_type=preprocess_type
            )
            preprocess_ms = (time.monotonic() - t0) * 1000
            logger.info(f"🧹 [2/5] Preprocessing done: {len(text):,} chars in {preprocess_ms:.0f}ms")

            # Step 3: Chunk'la
            t0 = time.monotonic()
            logger.info("✂️ [3/5] Intelligent chunking...")

            chunker = IntelligentChunker(chunk_size=1000, chunk_overlap=200)
            chunk_objects = chunker.chunk(text)
            chunks = [c['content'] for c in chunk_objects]
            doc_label = filename.rsplit(".", 1)[0].replace("-", " ").replace("_", " ").title()
            doc_type = self._detect_document_type(filename, text)
            is_dictionary = self._is_dictionary_doc(filename, text)

            # Sayfa haritası: metindeki ## Sayfa N konumlarından chunk → sayfa eşlemesi
            _page_positions = [
                (m.start(), int(m.group(1)))
                for m in re.finditer(r"(?m)^#+\s*Sayfa\s+(\d+)\s*$", text)
            ]

            def _find_page(content: str) -> Optional[str]:
                if not _page_positions:
                    return None
                snippet = content[:200].strip()
                if not snippet:
                    return None
                pos = text.find(snippet)
                if pos == -1:
                    snippet = content[:80].strip()
                    pos = text.find(snippet)
                if pos == -1:
                    return None
                page = None
                for pp_start, pp_num in _page_positions:
                    if pp_start <= pos:
                        page = str(pp_num)
                    else:
                        break
                return page

            enriched: list[str] = []
            _chunk_pages: list[Optional[str]] = []
            for idx, c in enumerate(chunks):
                header = chunk_objects[idx].get("header", "")
                source_page = _find_page(c)
                _chunk_pages.append(source_page)
                prefix = f"[Kaynak: {doc_label} | Doküman Türü: {doc_type}]"
                if header:
                    prefix += f"\n[Bölüm: {header}]"
                enriched.append(f"{prefix}\n\n{c}")
            chunks = enriched
            chunk_ms = (time.monotonic() - t0) * 1000
            logger.info(
                f"✂️ [3/5] Chunking done: {len(chunks)} chunks in {chunk_ms:.0f}ms "
                f"(type={doc_type}, dictionary={is_dictionary})"
            )

            # Step 4: Embedding'ler oluştur
            t0 = time.monotonic()
            logger.info(f"📊 [4/5] Generating embeddings for {len(chunks)} chunks...")

            embeddings = await self.embedding_service.embed_batch(chunks)

            if len(embeddings) != len(chunks):
                raise Exception(f"Embedding count mismatch: {len(embeddings)} vs {len(chunks)}")

            embed_ms = (time.monotonic() - t0) * 1000
            logger.info(f"📊 [4/5] Embeddings done: {len(embeddings)} vectors in {embed_ms:.0f}ms")

            # Step 5: Veritabanına kaydet
            t0 = time.monotonic()
            logger.info(f"💾 [5/5] Saving {len(chunks)} chunks to database...")

            sections = (self._structured_json or {}).get("sections") or {}
            bi_list = sections.get("banka_istihbarati")
            piyasa_dict = sections.get("piyasa_istihbarati")
            if isinstance(piyasa_dict, dict) and "source_pages" in piyasa_dict:
                piyasa_copy = {k: v for k, v in piyasa_dict.items() if k != "source_pages"}
                piyasa_dict = piyasa_copy if piyasa_copy else None

            documents = []
            for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                meta = {
                    "file_type": Path(filename).suffix,
                    "chunk_size": len(chunk),
                    "original_length": original_length,
                    "header": chunk_objects[idx].get("header"),
                    "source": filename,
                    "chunk_position": chunk_objects[idx].get("position"),
                    "doc_type": doc_type,
                    "is_dictionary": is_dictionary,
                }
                sp = _chunk_pages[idx] if idx < len(_chunk_pages) else None
                if sp:
                    meta["source_page"] = sp
                if idx == 0:
                    if bi_list is not None:
                        meta["banka_istihbarati"] = bi_list
                    if piyasa_dict is not None:
                        meta["piyasa_istihbarati"] = piyasa_dict
                documents.append(
                    DocumentChunk(
                        filename=filename,
                        chunk_index=idx,
                        content=chunk,
                        embedding=embedding,
                        metadata=meta,
                        collection=collection,
                    )
                )

            saved_docs = await self.document_repository.save_batch(documents)
            save_ms = (time.monotonic() - t0) * 1000
            total_tokens = original_length // 4
            pipeline_ms = (time.monotonic() - pipeline_start) * 1000

            logger.info(
                f"✅ Ingestion completed: {filename} | "
                f"{len(saved_docs)} chunks, ~{total_tokens:,} tokens | "
                f"total={pipeline_ms:.0f}ms "
                f"(extract={extract_ms:.0f}, preprocess={preprocess_ms:.0f}, "
                f"chunk={chunk_ms:.0f}, embed={embed_ms:.0f}, save={save_ms:.0f})"
            )
 
            return DocumentIngestionResult(

                status="success",

                filename=filename,

                chunks_ingested=len(saved_docs),

                total_tokens=total_tokens,

                timestamp=datetime.now()

            )
 
        except Exception as e:

            logger.error(f"❌ Document ingestion failed: {str(e)}")

            return DocumentIngestionResult(

                status="error",

                filename=filename,

                chunks_ingested=0,

                total_tokens=0,

                timestamp=datetime.now(),

                error=str(e)

            )
 
